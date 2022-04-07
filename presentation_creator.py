from pptx import Presentation
from pptx.util import Inches
import googlesearch, requests, re
from bs4 import BeautifulSoup
from duckduckgo_images_api import search as ddgo
import time
inp = input()



search = list(googlesearch.search(inp, lang="ru", num_results=10))

for elem in search:
    if 'wikipedia' in elem:
        req = requests.get(elem).text

soup = BeautifulSoup(req, 'html.parser')
title1 = soup.select('h1')[0].text

subtitle = soup.find(class_='mw-parser-output').find('p').text
subtitle = re.findall(r'[А-я .,ё-]',subtitle)
subtitle = ''.join(subtitle)

##PRESENTATION
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
lyt = prs.slide_layouts[0]
slide = prs.slides.add_slide(lyt)
tit = slide.shapes.title
subtit = slide.placeholders[1]
tit.text = title1
subtit.text = subtitle





exep = 0
for i in range(10):
    try:
        time.sleep(1)
        main = soup.findAll('h2')[i]

        title = main.find_next(class_='mw-headline').text
        title = re.findall(r'[А-я .,-]', title)
        title = ''.join(title)

        subtitle = main.find_next('p').text
        subtitle = re.findall(r'[А-я .,-ё]', subtitle)
        subtitle = ''.join(subtitle)

        lyt = prs.slide_layouts[0]
        slide = prs.slides.add_slide(lyt)
        tit = slide.shapes.title
        subtit = slide.placeholders[1]
        tit.text = title
        subtit.text = subtitle

        img_path = 'presentatino/image'+ str(i)
        try:
            results = ddgo(title1, max_results=2)

            response = requests.get(results[i]['image'])

            file = open(img_path, "wb")
            file.write(response.content)

            slide.shapes.add_picture(img_path, Inches(0), Inches(0))
        except:
            pass

    except Exception as e:
        print(e)
        if exep > 2:
            break
        exep+=1
        continue



prs.save(title1+'.pptx')
